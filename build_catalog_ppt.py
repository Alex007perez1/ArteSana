from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from PIL import Image


BASE = Path(__file__).parent
OUT = BASE / "catalogo-artesana.pptx"

COLORS = {
    "hueso": RGBColor(251, 245, 234),
    "verde": RGBColor(85, 107, 47),
    "verde_profundo": RGBColor(45, 90, 39),
    "guindo": RGBColor(115, 33, 53),
    "rosa": RGBColor(251, 227, 228),
    "terracota": RGBColor(183, 101, 69),
    "gris": RGBColor(47, 46, 43),
    "blanco": RGBColor(255, 255, 255),
}


def rgb(name):
    return COLORS[name]


def add_rect(slide, x, y, w, h, color, transparency=0, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.fill.transparency = transparency
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    return shp


def add_text(slide, text, x, y, w, h, size=18, color=None, bold=False, font="Montserrat", align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or rgb("gris")
    return box


def add_multiline(slide, lines, x, y, w, h, size=12.5, color=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        p.level = 0
        run = p.add_run()
        run.text = line
        run.font.name = "Montserrat"
        run.font.size = Pt(size)
        run.font.color.rgb = color or rgb("gris")
    return box


def add_picture_fit(slide, img_name, x, y, w, h):
    path = BASE / img_name
    with Image.open(path) as im:
        iw, ih = im.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        pic_w = w
        pic_h = w / img_ratio
    else:
        pic_h = h
        pic_w = h * img_ratio
    return slide.shapes.add_picture(str(path), Inches(x + (w - pic_w) / 2), Inches(y + (h - pic_h) / 2), width=Inches(pic_w), height=Inches(pic_h))


def add_frame(slide):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.28), Inches(0.28), Inches(12.77), Inches(6.94))
    shp.fill.background()
    shp.line.color.rgb = RGBColor(204, 190, 168)
    shp.line.width = Pt(0.8)


def add_pill(slide, text, x, y, w, color, text_color=None):
    add_rect(slide, x, y, w, 0.38, color, radius=True)
    add_text(slide, text, x + 0.05, y + 0.07, w - 0.1, 0.22, 11, text_color or rgb("guindo"), True, align=PP_ALIGN.CENTER)


def product_slide(prs, data, reverse=False, accent="rosa"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.333, 7.5, rgb("hueso"))
    if reverse:
        content_x, visual_x = 0.62, 7.15
    else:
        visual_x, content_x = 0.0, 6.15

    add_rect(slide, visual_x, 0, 6.0, 7.5, rgb(accent))
    add_rect(slide, visual_x + 0.76, 0.95, 4.5, 4.85, rgb("blanco"), transparency=12, radius=True)
    add_picture_fit(slide, data["image"], visual_x + 0.72, 0.68, 4.6, 5.58)

    add_picture_fit(slide, "1_20260517_195608_0000.png", 11.52, 6.48, 1.05, 0.55)
    add_text(slide, data["category"].upper(), content_x, 0.78, 5.75, 0.32, 9.5, rgb("terracota"), True)
    add_text(slide, data["name"], content_x, 1.05, 5.85, 0.95, 36, rgb("guindo"), False, "Georgia")
    add_text(slide, data["intro"], content_x, 1.98, 5.78, 0.62, 15, rgb("gris"))
    add_pill(slide, data["presentation"], content_x, 2.72, 1.55, rgb("rosa"))
    add_pill(slide, "Bs ___", content_x + 1.72, 2.72, 1.2, rgb("guindo"), rgb("blanco"))

    add_rect(slide, content_x, 3.28, 2.72, 1.08, rgb("blanco"), transparency=8, line=RGBColor(225, 207, 194), radius=True)
    add_text(slide, "BENEFICIO", content_x + 0.16, 3.42, 2.38, 0.22, 8.5, rgb("verde_profundo"), True)
    add_text(slide, data["benefit"], content_x + 0.16, 3.66, 2.38, 0.5, 10.5, rgb("gris"))

    add_rect(slide, content_x + 3.0, 3.28, 2.72, 1.08, rgb("blanco"), transparency=8, line=RGBColor(225, 207, 194), radius=True)
    add_text(slide, "USO", content_x + 3.16, 3.42, 2.38, 0.22, 8.5, rgb("verde_profundo"), True)
    add_text(slide, data["use"], content_x + 3.16, 3.66, 2.38, 0.5, 10.5, rgb("gris"))

    add_rect(slide, content_x, 4.58, 2.72, 1.08, rgb("blanco"), transparency=8, line=RGBColor(225, 207, 194), radius=True)
    add_text(slide, "APLICACION", content_x + 0.16, 4.72, 2.38, 0.22, 8.5, rgb("verde_profundo"), True)
    add_text(slide, data["apply"], content_x + 0.16, 4.96, 2.38, 0.5, 10.5, rgb("gris"))

    add_rect(slide, content_x + 3.0, 4.58, 2.72, 1.08, rgb("blanco"), transparency=8, line=RGBColor(225, 207, 194), radius=True)
    add_text(slide, "RECOMENDACION", content_x + 3.16, 4.72, 2.38, 0.22, 8.5, rgb("verde_profundo"), True)
    add_text(slide, data["recommend"], content_x + 3.16, 4.96, 2.38, 0.5, 10.5, rgb("gris"))

    add_rect(slide, content_x, 5.96, 5.72, 0.55, rgb("rosa"), transparency=0, radius=True)
    add_text(slide, data["note"], content_x + 0.18, 6.12, 5.35, 0.25, 10.5, rgb("guindo"), True)
    add_frame(slide)


products = [
    {
        "name": "Jabon de Carbon",
        "category": "Jabones artesanales",
        "image": "1779057926520.jpg",
        "presentation": "90 g",
        "intro": "Limpieza profunda para piel con exceso de grasa, brillo e impurezas visibles.",
        "benefit": "El carbon activado ayuda a adsorber sebo, particulas e impurezas superficiales.",
        "use": "Humedecer, espumar 30 a 45 segundos y retirar con agua.",
        "apply": "Usar 1 vez al dia en piel grasa. Alternar si reseca.",
        "recommend": "Aplicar hidratante despues para conservar la barrera cutanea.",
        "note": "Uso seguro siguiendo recomendaciones de aplicacion.",
        "accent": "verde",
    },
    {
        "name": "Jabon de Manzanilla",
        "category": "Piel sensible",
        "image": "1779058143113.jpg",
        "presentation": "90 g",
        "intro": "Suavidad herbal para limpiar, nutrir y dar confort a pieles delicadas.",
        "benefit": "La manzanilla aporta compuestos calmantes asociados al confort de la piel sensible.",
        "use": "Masajear con movimientos circulares suaves, sin friccion intensa.",
        "apply": "Ideal para uso diario en rostro, cuello o zonas delicadas.",
        "recommend": "Secar con toalla limpia a toques y complementar con crema ligera.",
        "note": "Formula artesanal para una limpieza amable y sensorial.",
        "accent": "rosa",
    },
    {
        "name": "Jabon de Romero",
        "category": "Herbal revitalizante",
        "image": "1779058219869.jpg",
        "presentation": "90 g",
        "intro": "Limpieza estimulante con aroma herbal para una piel fresca y tonificada.",
        "benefit": "El romero contiene antioxidantes naturales usados en cuidado cutaneo.",
        "use": "Crear espuma en manos y aplicar sobre piel humeda durante 30 segundos.",
        "apply": "Recomendado para ducha matinal o despues de actividad fisica ligera.",
        "recommend": "Guardar en jabonera drenante para conservar su textura.",
        "note": "Cuidado artesanal con sensacion herbal y vigorizante.",
        "accent": "terracota",
    },
    {
        "name": "Jabon de Aloe Vera",
        "category": "Humectacion diaria",
        "image": "1779059711115.jpg",
        "presentation": "90 g",
        "intro": "Una barra suave para acompanar la hidratacion y el cuidado cotidiano de la piel.",
        "benefit": "El aloe vera es valorado por mucilagos con efecto humectante.",
        "use": "Aplicar sobre piel humeda, espumar suavemente y enjuagar.",
        "apply": "Puede usarse en rostro y cuerpo como limpieza diaria suave.",
        "recommend": "Combinar con crema facial para mayor suavidad al tacto.",
        "note": "Uso seguro siguiendo recomendaciones de aplicacion.",
        "accent": "verde",
    },
    {
        "name": "Jabon de Arroz",
        "category": "Luminosidad natural",
        "image": "1779059812483.jpg",
        "presentation": "90 g",
        "intro": "Para una piel mas suave, uniforme y luminosa dentro de una rutina constante.",
        "benefit": "El arroz aporta almidones y compuestos suavizantes usados en cosmetica natural.",
        "use": "Masajear la espuma evitando el contorno de ojos.",
        "apply": "Usar en la noche y complementar de dia con hidratacion.",
        "recommend": "Combo sugerido: Jabon de Arroz + Crema Facial + Shampoo.",
        "note": "Ideal para suavidad, tono visual mas parejo y cuidado constante.",
        "accent": "rosa",
    },
    {
        "name": "Jabon de Azufre",
        "category": "Control de brillo",
        "image": "1779059889457.jpg",
        "presentation": "90 g",
        "intro": "Limpieza especifica para piel grasa, brillo persistente y poros congestionados.",
        "benefit": "El azufre se usa en cosmetica por su apoyo en higiene de piel grasa.",
        "use": "Aplicar solo en zonas necesarias y enjuagar bien.",
        "apply": "Iniciar 3 veces por semana. Aumentar solo si se tolera.",
        "recommend": "No combinar el mismo dia con exfoliacion intensa.",
        "note": "Uso seguro siguiendo recomendaciones de aplicacion.",
        "accent": "terracota",
    },
    {
        "name": "Crema Facial",
        "category": "Hidratacion nutritiva",
        "image": "1779058319163.jpg",
        "presentation": "Gramos",
        "intro": "Sabila, colageno, pepino y vitaminas A, E y C para piel nutrida y flexible.",
        "benefit": "Apoya la hidratacion superficial y la sensacion de elasticidad.",
        "use": "Aplicar una pequena cantidad sobre rostro limpio, manana y noche.",
        "apply": "Distribuir hacia arriba en mejillas, frente, cuello y escote.",
        "recommend": "Usar despues del Jabon de Arroz para una rutina de luminosidad.",
        "note": "Textura para cuidado diario con enfoque artesanal natural.",
        "accent": "verde",
    },
    {
        "name": "Crema Anti-Stress",
        "category": "Cuidado corporal",
        "image": "1779058568481.jpg",
        "presentation": "150 g",
        "intro": "Formula relajante con magnesio y aceites esenciales para masaje localizado.",
        "benefit": "El masaje favorece la sensacion de descanso muscular y bienestar corporal.",
        "use": "Aplicar en cuello, hombros, espalda o piernas con movimientos lentos.",
        "apply": "Usar por la noche o despues de esfuerzo fisico.",
        "recommend": "Masajear 3 a 5 minutos hasta absorcion confortable.",
        "note": "Uso externo. Uso seguro siguiendo recomendaciones.",
        "accent": "rosa",
    },
    {
        "name": "Crema Vital-Flex",
        "category": "Alivio corporal rapido",
        "image": "1779058638197.jpg",
        "presentation": "150 g",
        "intro": "Crema para masaje en zonas de tension con sensacion de calor intenso.",
        "benefit": "El masaje localizado ayuda a activar la sensacion termica y confort.",
        "use": "Aplicar poca cantidad sobre la zona y masajear hasta absorcion.",
        "apply": "Evitar ojos, mucosas o piel recien rasurada.",
        "recommend": "Lavar manos despues y no cubrir con calor adicional.",
        "note": "Uso seguro siguiendo recomendaciones de aplicacion.",
        "accent": "terracota",
    },
    {
        "name": "Shampoo Anti-Stress",
        "category": "Cuidado capilar",
        "image": "1779058854107.jpg",
        "presentation": "160 g",
        "intro": "Refrescante y relajante, con aceites esenciales y extractos botanicos.",
        "benefit": "Apoya la limpieza capilar con sensacion fresca, suave y aromatica.",
        "use": "Aplicar sobre cabello mojado, masajear cuero cabelludo y enjuagar.",
        "apply": "Repetir si se requiere mayor limpieza. Usar 2 a 4 veces por semana.",
        "recommend": "Combo especial con Jabon de Arroz y Crema Facial.",
        "note": "Producto unisex para una rutina de frescura natural.",
        "accent": "verde",
    },
    {
        "name": "Kefir de Leche",
        "category": "Probiotico artesanal",
        "image": "image.jpg",
        "presentation": "1 L / 1/2 L",
        "intro": "Bebida probiotica artesanal de consumo diario para habitos de bienestar.",
        "benefit": "Contiene microorganismos asociados al equilibrio de la microbiota intestinal.",
        "use": "Consumir frio. Agitar suavemente antes de servir.",
        "apply": "Iniciar con porciones pequenas y aumentar segun preferencia.",
        "recommend": "Mantener refrigerado y consumir dentro del periodo indicado.",
        "note": "Sin azucar, sin conservantes, sin lactosa y probiotico artesanal.",
        "accent": "rosa",
    },
]


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, rgb("rosa"))
add_rect(slide, 7.25, 0, 6.1, 7.5, rgb("verde"))
add_picture_fit(slide, "1_20260517_195608_0000.png", 0.9, 0.65, 4.55, 2.55)
add_text(slide, "CATALOGO REVISTA 2026", 0.95, 3.25, 4.4, 0.28, 10, rgb("terracota"), True)
add_text(slide, "Bienestar natural para cada dia", 0.9, 3.58, 5.6, 1.35, 38, rgb("guindo"), False, "Georgia")
add_text(slide, "Productos artesanales seleccionados para acompanar el cuidado de la piel, el cabello y el equilibrio diario.", 0.95, 5.18, 5.2, 0.72, 16, rgb("gris"))
add_pill(slide, "CREANDO BIENESTAR INTEGRAL", 0.95, 6.25, 3.25, rgb("guindo"), rgb("blanco"))
add_picture_fit(slide, "image.jpg", 7.75, 0.58, 4.82, 6.4)
add_frame(slide)

for i, data in enumerate(products):
    product_slide(prs, data, reverse=bool(i % 2), accent=data["accent"])

# Contact slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, rgb("hueso"))
add_rect(slide, 0, 0, 5.55, 7.5, rgb("rosa"))
add_text(slide, "HAZ TU PEDIDO", 0.78, 0.82, 4.2, 0.28, 10, rgb("terracota"), True)
add_text(slide, "ArteSana", 0.78, 1.15, 4.2, 0.75, 36, rgb("guindo"), False, "Georgia")
add_text(slide, "Productos artesanales para una rutina natural, practica y sensorial.", 0.82, 1.98, 4.15, 0.56, 15.5, rgb("gris"))
add_picture_fit(slide, "1_20260517_195608_0000.png", 0.9, 2.85, 3.45, 1.9)
add_text(slide, "Contactanos al +591 68703773", 0.82, 5.24, 4.3, 0.35, 17, rgb("guindo"), True)
add_text(slide, "Facebook: ArteSana", 0.82, 5.7, 4.3, 0.28, 13, rgb("gris"))
add_text(slide, "Delivery gratis hasta 4to anillo - Santa Cruz", 0.82, 6.05, 4.4, 0.32, 13, rgb("gris"))

add_rect(slide, 6.08, 0.82, 2.82, 3.52, rgb("blanco"), transparency=5, line=RGBColor(225, 207, 194), radius=True)
add_picture_fit(slide, "3_20260517_195608_0002.png", 6.34, 1.05, 2.3, 2.3)
add_text(slide, "WhatsApp", 6.24, 3.5, 2.5, 0.33, 18, rgb("guindo"), False, "Georgia", PP_ALIGN.CENTER)
add_text(slide, "Escanea y realiza tu pedido.", 6.22, 3.88, 2.5, 0.24, 10.5, rgb("gris"), align=PP_ALIGN.CENTER)

add_rect(slide, 9.48, 0.82, 2.82, 3.52, rgb("blanco"), transparency=5, line=RGBColor(225, 207, 194), radius=True)
add_picture_fit(slide, "2_20260517_195608_0001.png", 9.74, 1.05, 2.3, 2.3)
add_text(slide, "Pago digital", 9.62, 3.5, 2.5, 0.33, 18, rgb("guindo"), False, "Georgia", PP_ALIGN.CENTER)
add_text(slide, "Forma de pago digital.", 9.64, 3.88, 2.5, 0.24, 10.5, rgb("gris"), align=PP_ALIGN.CENTER)

add_rect(slide, 6.08, 4.78, 6.22, 1.32, rgb("guindo"), radius=True)
add_text(slide, "Combos especiales", 6.4, 4.99, 3.7, 0.35, 22, rgb("blanco"), False, "Georgia")
add_text(slide, "Jabon de Arroz + Crema Facial + Shampoo Anti-Stress. Pregunta por disponibilidad y precio final.", 6.4, 5.42, 4.4, 0.38, 11.5, rgb("blanco"))
add_pill(slide, "Bs ___", 10.9, 5.2, 0.95, rgb("rosa"), rgb("guindo"))
add_frame(slide)

prs.save(OUT)
print(OUT)
